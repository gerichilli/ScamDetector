from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "ScamDetector_Diagram_Pack.pptx"

W, H = Inches(13.333), Inches(7.5)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
PAPER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(37, 99, 235)
BLUE_DARK = RGBColor(30, 64, 175)
RED = RGBColor(185, 28, 28)
RED_PALE = RGBColor(254, 242, 242)
AMBER = RGBColor(180, 83, 9)
AMBER_PALE = RGBColor(255, 251, 235)
GREEN = RGBColor(4, 120, 87)
GREEN_PALE = RGBColor(236, 253, 245)
LINE = RGBColor(203, 213, 225)
SLATE = RGBColor(226, 232, 240)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color=LINE, width=1.4):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, section, title, subtitle=None):
    add_text(slide, section, 0.55, 0.35, 0.65, 0.34, 13, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, title, 1.25, 0.28, 10.4, 0.62, 25, INK, True)
    if subtitle:
        add_text(slide, subtitle, 1.25, 0.91, 10.6, 0.36, 11.5, MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.02))
    set_fill(line, SLATE)
    line.line.fill.background()


def add_footer(slide, n):
    add_text(slide, "ScamDetector · Diagram pack", 0.55, 7.08, 4, 0.25, 9.5, MUTED)
    add_text(slide, f"{n:02d}", 12.15, 7.02, 0.55, 0.28, 10, MUTED, True, PP_ALIGN.RIGHT)


def node(slide, text, x, y, w, h, fill=WHITE, border=LINE, color=INK, size=13, bold=True, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shp, fill)
    set_line(shp, border, 1.5)
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shp


def label(slide, text, x, y, w, h, color=MUTED, size=10):
    return add_text(slide, text, x, y, w, h, size, color, False, PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color=RGBColor(100, 116, 139), width=1.7):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    try:
        c.line.end_arrowhead = True
    except Exception:
        pass
    return c


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    return slide


def cover(prs):
    s = blank(prs)
    add_text(s, "BÁO CÁO ĐỒ ÁN · SE104", 0.75, 0.72, 4.6, 0.35, 12, BLUE, True)
    add_text(s, "ScamDetector", 0.75, 1.35, 6.8, 0.75, 42, INK, True)
    add_text(s, "Diagram Pack", 0.78, 2.07, 4, 0.45, 24, RED, True)
    add_text(s, "Bộ slide tập trung vào sơ đồ: DFD, Use Case, Activity, Sequence, ERD, Security, Deployment và AI/OCR pipeline.", 0.8, 2.83, 5.4, 1.1, 16, MUTED)
    node(s, "Frontend\nReact + Vite", 7.2, 1.0, 1.7, 0.9, WHITE, LINE, BLUE_DARK)
    node(s, "Backend\nFastAPI", 9.45, 1.0, 1.7, 0.9, WHITE, LINE, BLUE_DARK)
    node(s, "PostgreSQL\nScam DB", 9.45, 2.55, 1.7, 0.9, WHITE, LINE, GREEN)
    node(s, "AI Model\nMasked input", 7.2, 2.55, 1.7, 0.9, RED_PALE, RGBColor(252, 165, 165), RED)
    arrow(s, 8.9, 1.45, 9.45, 1.45)
    arrow(s, 10.3, 1.9, 10.3, 2.55)
    arrow(s, 9.45, 3.0, 8.9, 3.0)
    arrow(s, 8.05, 2.55, 8.05, 1.9)
    add_footer(s, 1)


def agenda(prs):
    s = blank(prs)
    add_title(s, "00", "Diagram map theo deck mẫu điểm cao", "Mỗi loại sơ đồ có phiên bản tương ứng cho hệ thống ScamDetector.")
    items = [
        ("DFD", "Context + Level 1 cho tra cứu, báo cáo, chatbot"),
        ("Use Case", "User, moderator/admin, AI/OCR service"),
        ("Activity", "Chat AI và OCR ảnh"),
        ("Sequence", "Backend gọi AI model sau bước masking"),
        ("Data Design", "ERD + module/class view"),
        ("Ops & Security", "Docker deployment + privacy boundary"),
    ]
    for i, (a, b) in enumerate(items):
        x = 0.8 + (i % 2) * 6.0
        y = 1.7 + (i // 2) * 1.35
        node(s, a, x, y, 1.35, 0.68, BLUE, BLUE, WHITE, 15)
        add_text(s, b, x + 1.55, y + 0.1, 4.1, 0.45, 14, INK, True)
    add_footer(s, 2)


def dfd0(prs):
    s = blank(prs)
    add_title(s, "01", "DFD mức 0: ScamDetector là điểm trung gian cảnh báo", "Người dùng không gửi trực tiếp dữ liệu nhạy cảm đến AI; backend kiểm soát luồng dữ liệu.")
    node(s, "Người dùng\nngười cao tuổi / thân nhân", 0.8, 3.0, 2.2, 0.95, WHITE, LINE)
    node(s, "ScamDetector\nweb platform", 5.1, 2.65, 2.4, 1.25, BLUE, BLUE, WHITE, 17)
    node(s, "CSDL cảnh báo\nsố điện thoại / báo cáo", 9.8, 1.7, 2.3, 0.95, WHITE, LINE, GREEN)
    node(s, "AI provider\nOpenAI API", 9.8, 3.6, 2.3, 0.95, RED_PALE, RGBColor(252, 165, 165), RED)
    node(s, "Admin / Moderator\nxác minh báo cáo", 9.8, 5.35, 2.3, 0.95, WHITE, LINE)
    arrow(s, 3.0, 3.45, 5.1, 3.25)
    label(s, "tin nhắn / ảnh / số cần kiểm tra", 3.2, 2.95, 1.8, 0.25)
    arrow(s, 7.5, 2.95, 9.8, 2.15)
    arrow(s, 7.5, 3.25, 9.8, 4.05)
    arrow(s, 7.5, 3.55, 9.8, 5.8)
    arrow(s, 5.1, 3.75, 3.0, 3.75)
    label(s, "verdict: safe / warning / danger", 3.35, 3.95, 1.7, 0.3, RED)
    add_footer(s, 3)


def dfd1(prs):
    s = blank(prs)
    add_title(s, "02", "DFD mức 1: bốn luồng dữ liệu chính", "Tách rõ luồng tra cứu, báo cáo, thống kê và chatbot AI/OCR.")
    centers = [
        ("1. Tra cứu", "phone → normalize → match DB", 0.75, 1.65),
        ("2. Báo cáo", "report → evidence → approval", 4.05, 1.65),
        ("3. Thống kê", "alerts → aggregate → chart", 7.35, 1.65),
        ("4. Chatbot AI", "text/image → mask → model", 10.65, 1.65),
    ]
    for title, desc, x, y in centers:
        node(s, title, x, y, 2.1, 0.65, BLUE, BLUE, WHITE, 13)
        node(s, desc, x, y + 0.82, 2.1, 0.78, WHITE, LINE, INK, 10, False)
    node(s, "Frontend", 1.0, 5.0, 1.8, 0.62, WHITE, LINE)
    node(s, "API Router", 3.65, 5.0, 1.8, 0.62, WHITE, LINE)
    node(s, "Services", 6.3, 5.0, 1.8, 0.62, WHITE, LINE)
    node(s, "Database / AI", 8.95, 5.0, 2.0, 0.62, WHITE, LINE)
    for x in [2.8, 5.45, 8.1]:
        arrow(s, x, 5.31, x + 0.85, 5.31)
    add_footer(s, 4)


def use_case(prs):
    s = blank(prs)
    add_title(s, "03", "Use Case diagram: vai trò và quyền thao tác", "Tập trung vào trải nghiệm cảnh báo, báo cáo và quản trị dữ liệu lừa đảo.")
    node(s, "USER", 0.75, 2.6, 1.4, 0.7, BLUE, BLUE, WHITE)
    node(s, "ADMIN /\nMODERATOR", 0.75, 4.3, 1.4, 0.8, INK, INK, WHITE, 12)
    boundary = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(1.45), Inches(8.7), Inches(4.8))
    set_fill(boundary, RGBColor(241, 245, 249))
    set_line(boundary, LINE, 1.2)
    add_text(s, "ScamDetector System Boundary", 2.95, 1.62, 3.4, 0.3, 11, MUTED, True)
    cases = [
        ("Tra cứu số", 3.15, 2.25), ("Chat cảnh báo", 5.25, 2.25), ("Đọc ảnh OCR", 7.35, 2.25),
        ("Gửi báo cáo", 3.15, 3.55), ("Xem thống kê", 5.25, 3.55), ("Duyệt báo cáo", 7.35, 3.55),
        ("Quản lý DB scam", 4.2, 4.85), ("Xem lịch sử", 6.55, 4.85),
    ]
    for t, x, y in cases:
        shp = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(1.55), Inches(0.62))
        set_fill(shp, WHITE)
        set_line(shp, RGBColor(148, 163, 184), 1.1)
        shp.text_frame.text = t
        p = shp.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = "Arial"
        p.runs[0].font.size = Pt(10.5)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = INK
    for x, y in [(3.15,2.56),(5.25,2.56),(7.35,2.56),(3.15,3.86),(5.25,3.86),(6.55,5.16)]:
        arrow(s, 2.15, 2.95, x, y, LINE, 1.0)
    for x, y in [(7.35,3.86),(4.2,5.16),(5.25,3.86)]:
        arrow(s, 2.15, 4.7, x, y, LINE, 1.0)
    add_footer(s, 5)


def activity_chat(prs):
    s = blank(prs)
    add_title(s, "04", "Activity diagram: luồng chat AI cảnh báo lừa đảo", "Luồng này minh họa rõ điểm bảo mật: masking trước khi gọi model.")
    steps = [
        ("Nhập câu hỏi", BLUE, WHITE),
        ("Backend nhận request", WHITE, INK),
        ("Mask dữ liệu nhạy cảm", RED_PALE, RED),
        ("Gọi AI model", WHITE, INK),
        ("Chuẩn hóa verdict", WHITE, INK),
        ("Hiển thị bubble màu", GREEN_PALE, GREEN),
    ]
    x = 1.0
    for i, (t, fill, color) in enumerate(steps):
        node(s, t, x + i * 1.95, 3.0, 1.55, 0.75, fill, LINE if fill == WHITE else fill, color, 11)
        if i < len(steps) - 1:
            arrow(s, x + i * 1.95 + 1.55, 3.38, x + (i + 1) * 1.95, 3.38)
    add_text(s, "Fallback: nếu API lỗi hoặc key chưa hợp lệ → rule-based detector vẫn trả lời để UI không hỏng.", 1.25, 5.0, 10.8, 0.45, 15, MUTED, True, PP_ALIGN.CENTER)
    add_footer(s, 6)


def activity_ocr(prs):
    s = blank(prs)
    add_title(s, "05", "Activity diagram: đọc ảnh rồi phân tích nội dung", "OCR chạy ở frontend; chỉ đoạn text trích xuất mới đi vào pipeline phân tích.")
    node(s, "Upload ảnh", 0.85, 2.0, 1.55, 0.7, BLUE, BLUE, WHITE)
    node(s, "Tesseract.js\nOCR trong browser", 3.0, 2.0, 1.85, 0.9, WHITE, LINE)
    node(s, "Có text rõ?", 5.6, 1.95, 1.55, 0.82, AMBER_PALE, RGBColor(252, 211, 77), AMBER)
    node(s, "Gửi text đến backend", 8.0, 1.45, 1.85, 0.72, WHITE, LINE)
    node(s, "Hiện hướng dẫn\nnhập lại thủ công", 8.0, 2.7, 1.85, 0.82, WHITE, LINE)
    node(s, "AI / rule verdict", 10.65, 1.45, 1.65, 0.72, RED_PALE, RGBColor(252, 165, 165), RED)
    node(s, "Bubble cảnh báo", 10.65, 2.7, 1.65, 0.72, GREEN_PALE, RGBColor(110, 231, 183), GREEN)
    arrow(s, 2.4, 2.35, 3.0, 2.35)
    arrow(s, 4.85, 2.35, 5.6, 2.35)
    arrow(s, 7.15, 2.2, 8.0, 1.83)
    label(s, "yes", 7.2, 1.73, 0.45, 0.22, GREEN)
    arrow(s, 7.15, 2.45, 8.0, 3.1)
    label(s, "no", 7.2, 2.78, 0.45, 0.22, RED)
    arrow(s, 9.85, 1.82, 10.65, 1.82)
    arrow(s, 11.47, 2.17, 11.47, 2.7)
    add_footer(s, 7)


def sequence_ai(prs):
    s = blank(prs)
    add_title(s, "06", "Sequence diagram: backend gọi AI model an toàn", "API key chỉ nằm ở backend; frontend không bao giờ biết key.")
    actors = ["User", "ChatbotWidget", "FastAPI /chat", "Redaction", "OpenAI API", "Response UI"]
    xs = [0.8, 2.8, 4.95, 7.05, 9.15, 11.15]
    for x, a in zip(xs, actors):
        node(s, a, x, 1.55, 1.35, 0.55, BLUE if a == "OpenAI API" else WHITE, BLUE if a == "OpenAI API" else LINE, WHITE if a == "OpenAI API" else INK, 10.5)
        ln = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.67), Inches(2.15), Inches(0.015), Inches(3.9))
        set_fill(ln, LINE)
        ln.line.fill.background()
    messages = [
        (0, 1, 2.5, "submit(text)"),
        (1, 2, 3.0, "POST /chat/analyze"),
        (2, 3, 3.5, "redact_sensitive_text()"),
        (3, 4, 4.0, "masked prompt"),
        (4, 2, 4.55, "JSON verdict"),
        (2, 5, 5.1, "safe/warning/danger"),
    ]
    for a, b, y, txt in messages:
        arrow(s, xs[a] + 1.15, y, xs[b] + 0.2, y, RED if "masked" in txt else RGBColor(100,116,139), 1.4)
        label(s, txt, min(xs[a], xs[b]) + 0.55, y - 0.28, abs(xs[b]-xs[a]) + 0.9, 0.22, RED if "masked" in txt else MUTED, 9)
    add_footer(s, 8)


def sequence_report(prs):
    s = blank(prs)
    add_title(s, "07", "Sequence diagram: báo cáo số lừa đảo và duyệt dữ liệu", "Luồng này tạo nguồn dữ liệu nội bộ cho tra cứu và thống kê.")
    actors = ["User", "ReportPage", "Reports API", "ScamEntity", "AdminPage", "Stats API"]
    xs = [0.8, 2.85, 4.85, 6.85, 8.85, 10.85]
    for x, a in zip(xs, actors):
        node(s, a, x, 1.45, 1.4, 0.58, WHITE, LINE, INK, 10.5)
        ln = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.7), Inches(2.05), Inches(0.012), Inches(4.0))
        set_fill(ln, LINE)
        ln.line.fill.background()
    for a, b, y, txt in [
        (0, 1, 2.4, "submit report"),
        (1, 2, 2.9, "POST /reports"),
        (2, 3, 3.4, "upsert entity"),
        (4, 2, 4.15, "approve / reject"),
        (2, 5, 4.8, "aggregate alerts"),
        (5, 1, 5.35, "dashboard data"),
    ]:
        arrow(s, xs[a] + 1.12, y, xs[b] + 0.22, y)
        label(s, txt, min(xs[a], xs[b]) + 0.55, y - 0.28, abs(xs[b]-xs[a]) + 0.9, 0.22, MUTED, 9)
    add_footer(s, 9)


def erd(prs):
    s = blank(prs)
    add_title(s, "08", "ERD: dữ liệu cảnh báo xoay quanh scam entity", "Thiết kế hỗ trợ cả tra cứu nhanh, báo cáo cộng đồng và lịch sử cảnh báo.")
    tables = [
        ("users\nid PK\nemail\nrole", 0.8, 1.75),
        ("scam_entities\nid PK\nnormalized_value\nrisk_level", 3.2, 1.75),
        ("scam_reports\nid PK\nreporter_id FK\nscam_entity_id FK", 5.9, 1.75),
        ("call_records\nid PK\nelderly_user_id FK\nphone_number", 3.2, 4.3),
        ("scam_alerts\nid PK\ncall_id FK\nrisk_level", 5.9, 4.3),
        ("notifications\nid PK\nalert_id FK\ntarget_user_id FK", 8.6, 4.3),
        ("scam_database\nphone\npattern\nrisk_level", 8.6, 1.75),
    ]
    for t, x, y in tables:
        node(s, t, x, y, 2.0, 1.05, WHITE, LINE, INK, 10, False, False)
    for x1,y1,x2,y2 in [(2.8,2.25,3.2,2.25),(5.2,2.25,5.9,2.25),(4.2,2.8,4.2,4.3),(5.2,4.82,5.9,4.82),(7.9,4.82,8.6,4.82),(7.9,2.25,8.6,2.25)]:
        arrow(s,x1,y1,x2,y2,RGBColor(100,116,139),1.1)
    add_footer(s, 10)


def module_view(prs):
    s = blank(prs)
    add_title(s, "09", "Module/Class view: backend tách API, schema, service, model", "Cấu trúc này giúp thay rule-based bằng AI model mà không đổi giao diện.")
    layers = [
        ("Frontend", ["ChatbotWidget", "LookupPage", "StatsPage"], 0.9, BLUE),
        ("API Layer", ["chat.py", "lookup.py", "reports.py"], 3.5, GREEN),
        ("Service Layer", ["chatbot_service", "normalization", "alert_service"], 6.1, AMBER),
        ("Data Layer", ["User", "ScamReport", "ScamAlert"], 8.7, RED),
    ]
    for title, items, x, col in layers:
        node(s, title, x, 1.7, 2.1, 0.55, col, col, WHITE, 13)
        for i, item in enumerate(items):
            node(s, item, x, 2.55 + i * 0.8, 2.1, 0.55, WHITE, LINE, INK, 10, False)
    for x in [3.0,5.6,8.2]:
        arrow(s, x, 3.35, x+0.5, 3.35)
    add_text(s, "Key idea: UI chỉ cần biết endpoint; service quyết định dùng AI hay fallback rule-based.", 1.1, 5.75, 10.8, 0.4, 15, MUTED, True, PP_ALIGN.CENTER)
    add_footer(s, 11)


def security(prs):
    s = blank(prs)
    add_title(s, "10", "Security diagram: bảo vệ API key và dữ liệu nhạy cảm", "Ranh giới bảo mật nằm ở backend, trước khi gọi AI provider.")
    node(s, "Frontend\nkhông có API key", 0.95, 2.35, 2.0, 0.9, WHITE, LINE)
    node(s, "Backend\n.env + settings", 4.0, 2.35, 2.0, 0.9, BLUE, BLUE, WHITE)
    node(s, "Redaction\nPHONE / EMAIL / OTP", 7.05, 2.35, 2.1, 0.9, RED_PALE, RGBColor(252, 165, 165), RED)
    node(s, "AI Provider\nmasked prompt only", 10.25, 2.35, 2.0, 0.9, WHITE, LINE)
    for x in [2.95, 6.0, 9.15]:
        arrow(s, x, 2.8, x+1.05, 2.8)
    node(s, "Không log raw OTP / số tài khoản", 2.0, 4.45, 2.6, 0.65, GREEN_PALE, RGBColor(110, 231, 183), GREEN, 11)
    node(s, "Fallback nếu AI lỗi", 5.35, 4.45, 2.2, 0.65, AMBER_PALE, RGBColor(252, 211, 77), AMBER, 11)
    node(s, "Frontend chỉ nhận verdict", 8.3, 4.45, 2.6, 0.65, GREEN_PALE, RGBColor(110, 231, 183), GREEN, 11)
    add_footer(s, 12)


def deployment(prs):
    s = blank(prs)
    add_title(s, "11", "Deployment diagram: Docker Compose chạy 3 service chính", "Frontend, backend và database tách container, giao tiếp qua network nội bộ.")
    node(s, "Browser\nlocalhost:5173", 0.95, 2.8, 1.7, 0.75, WHITE, LINE)
    node(s, "frontend\nReact/Vite\nport 5173", 3.35, 2.55, 1.8, 1.0, BLUE, BLUE, WHITE)
    node(s, "backend\nFastAPI\nport 8000", 6.2, 2.55, 1.8, 1.0, GREEN, GREEN, WHITE)
    node(s, "db\nPostgres\nport 5432", 9.05, 2.55, 1.8, 1.0, WHITE, LINE, INK)
    node(s, "OpenAI API\ninternet", 6.2, 4.55, 1.8, 0.75, RED_PALE, RGBColor(252,165,165), RED)
    arrow(s, 2.65, 3.15, 3.35, 3.05)
    arrow(s, 5.15, 3.05, 6.2, 3.05)
    arrow(s, 8.0, 3.05, 9.05, 3.05)
    arrow(s, 7.1, 3.55, 7.1, 4.55)
    add_text(s, "docker-compose.yml: frontend depends_on backend; backend depends_on db; OPENAI_API_KEY nằm trong environment backend.", 1.1, 5.95, 11, 0.35, 13.5, MUTED, True, PP_ALIGN.CENTER)
    add_footer(s, 13)


def ai_pipeline(prs):
    s = blank(prs)
    add_title(s, "12", "AI pipeline: từ input đến màu cảnh báo trên UI", "Pipeline thống nhất cho cả text chat và text OCR từ ảnh.")
    stages = [
        ("Input\ntext / OCR", WHITE, INK),
        ("Normalize\ntrim + compact", WHITE, INK),
        ("Redact\nsensitive data", RED_PALE, RED),
        ("Model / Rule\nclassification", WHITE, INK),
        ("Verdict\nsafe / warning / danger", AMBER_PALE, AMBER),
        ("UI Bubble\ncolor-coded", GREEN_PALE, GREEN),
    ]
    for i, (t, fill, color) in enumerate(stages):
        node(s, t, 0.75 + i * 2.05, 2.2, 1.55, 0.9, fill, LINE if fill == WHITE else fill, color, 10.5)
        if i < len(stages)-1:
            arrow(s, 2.3 + i * 2.05, 2.65, 2.8 + i * 2.05, 2.65)
    node(s, "SAFE", 3.35, 4.55, 1.4, 0.6, GREEN_PALE, RGBColor(110,231,183), GREEN)
    node(s, "WARNING", 5.6, 4.55, 1.4, 0.6, AMBER_PALE, RGBColor(252,211,77), AMBER)
    node(s, "NGUY HIỂM", 7.85, 4.55, 1.55, 0.6, RED_PALE, RGBColor(252,165,165), RED)
    add_footer(s, 14)


def conclusion(prs):
    s = blank(prs)
    add_title(s, "13", "Slide diagram nào nên đưa vào deck chính?", "Chọn 6-8 slide mạnh nhất để thuyết trình; phần còn lại để appendix.")
    picks = [
        ("Bắt buộc", "DFD mức 0, DFD mức 1, Use Case, ERD"),
        ("Nên có", "Activity Chat, Sequence AI, Security diagram"),
        ("Appendix", "Module view, Deployment, OCR activity, Report sequence"),
    ]
    for i, (a,b) in enumerate(picks):
        node(s, a, 1.2, 1.8 + i*1.35, 1.8, 0.7, [BLUE,GREEN,AMBER][i], [BLUE,GREEN,AMBER][i], WHITE)
        add_text(s, b, 3.35, 1.9 + i*1.35, 7.9, 0.4, 17, INK, True)
    add_footer(s, 15)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    for fn in [
        cover, agenda, dfd0, dfd1, use_case, activity_chat, activity_ocr,
        sequence_ai, sequence_report, erd, module_view, security, deployment,
        ai_pipeline, conclusion,
    ]:
        fn(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
